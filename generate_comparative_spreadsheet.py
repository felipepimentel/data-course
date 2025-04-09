#!/usr/bin/env python3
import json
from collections import defaultdict
from pathlib import Path
from typing import Any, Dict, List, Optional, Set, Tuple

import numpy as np
import pandas as pd
from openpyxl import Workbook, load_workbook
from openpyxl.chart import LineChart, PieChart, RadarChart, Reference, Series
from openpyxl.formatting.rule import CellIsRule, ColorScaleRule
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter

# Importar constantes
from constants import CONCEPT_COLORS, FREQUENCY_LABELS, FREQUENCY_WEIGHTS

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print(
        "python-pptx não está instalado. A exportação para PowerPoint não estará disponível."
    )
    print("Para instalar: pip install python-pptx")


class ComparativeSpreadsheetGenerator:
    def __init__(self, base_path: str):
        self.base_path = Path(base_path)
        self.evaluations_by_person = defaultdict(dict)
        # Usar constantes importadas
        self.frequency_labels = FREQUENCY_LABELS
        self.frequency_weights = FREQUENCY_WEIGHTS
        # Cores para os conceitos
        self.concept_colors = CONCEPT_COLORS
        self.load_all_evaluations()
        # Track the criteria for each year
        self.year_criteria = self._extract_year_criteria()
        # Track all available years
        self.all_years = sorted(
            set(
                year
                for person_data in self.evaluations_by_person.values()
                for year in person_data.keys()
            )
        )
        # Track all available people
        self.all_people = sorted(self.evaluations_by_person.keys())

    def load_all_evaluations(self):
        """Load all evaluation files from the base path"""
        for person_dir in self.base_path.iterdir():
            if not person_dir.is_dir():
                continue

            person_name = person_dir.name

            for year_dir in person_dir.iterdir():
                if not year_dir.is_dir():
                    continue

                year = year_dir.name
                result_file = year_dir / "resultado.json"

                if result_file.exists():
                    with open(result_file, "r", encoding="utf-8") as f:
                        try:
                            data = json.load(f)
                            self.evaluations_by_person[person_name][year] = data
                        except json.JSONDecodeError:
                            print(f"Error decoding JSON from {result_file}")

    def _extract_year_criteria(self) -> Dict[str, Dict[str, Set[str]]]:
        """Extract all criteria (direcionadores and comportamentos) for each year"""
        year_criteria = defaultdict(lambda: defaultdict(set))

        # First pass: discover all unique years
        all_years = set()
        for person, years in self.evaluations_by_person.items():
            all_years.update(years.keys())

        # Second pass: extract criteria for each year
        for year in all_years:
            # Find a person who has data for this year
            for person, years in self.evaluations_by_person.items():
                if year in years and years[year]["success"] and "data" in years[year]:
                    data = years[year]["data"]

                    if "direcionadores" in data:
                        for direcionador in data["direcionadores"]:
                            dir_name = direcionador["direcionador"]

                            for comportamento in direcionador["comportamentos"]:
                                comp_name = comportamento["comportamento"]
                                year_criteria[year][dir_name].add(comp_name)

                    # One sample is enough
                    break

        return year_criteria

    def calculate_weighted_score(self, frequencies: List[int]) -> float:
        """Calculate a weighted score from frequency percentages"""
        if not frequencies or sum(frequencies) == 0:
            return 0

        # Calculate weighted sum using weights and frequencies
        weighted_sum = sum(
            freq * weight for freq, weight in zip(frequencies, self.frequency_weights)
        )

        # Normalize by total percentage (excluding n/a)
        total = sum(frequencies)

        return weighted_sum / total if total > 0 else 0

    def get_concept_by_year(self, person: str) -> Dict[str, str]:
        """Get the overall concept for a person across all available years"""
        result = {}
        for year, data in self.evaluations_by_person[person].items():
            if (
                data["success"]
                and "data" in data
                and "conceito_ciclo_filho_descricao" in data["data"]
            ):
                result[year] = data["data"]["conceito_ciclo_filho_descricao"]
            else:
                result[year] = "n/a"
        return result

    def get_behavior_scores(
        self, person: str, year: str
    ) -> Dict[str, Dict[str, Dict[str, Any]]]:
        """Extract behavior scores for all direcionadores"""
        result = {}
        if (
            person not in self.evaluations_by_person
            or year not in self.evaluations_by_person[person]
        ):
            return result

        data = self.evaluations_by_person[person][year]
        if (
            not data["success"]
            or "data" not in data
            or "direcionadores" not in data["data"]
        ):
            return result

        for direcionador in data["data"]["direcionadores"]:
            dir_name = direcionador["direcionador"]
            result[dir_name] = {}

            for comportamento in direcionador["comportamentos"]:
                comp_name = comportamento["comportamento"]
                scores = {}

                for avaliacao in comportamento["avaliacoes_grupo"]:
                    avaliador = avaliacao["avaliador"]
                    person_freq = avaliacao["frequencia_colaborador"]
                    group_freq = avaliacao["frequencia_grupo"]

                    person_score = self.calculate_weighted_score(person_freq)
                    group_score = self.calculate_weighted_score(group_freq)

                    scores[avaliador] = {
                        "freq_colaborador": person_freq,
                        "freq_grupo": group_freq,
                        "score_colaborador": person_score,
                        "score_grupo": group_score,
                        "difference": person_score - group_score,
                    }

                # Also get individual evaluations
                individual_evals = {}
                for avaliacao in comportamento["avaliacoes_individuais"]:
                    avaliador = avaliacao["avaliador"]
                    conceito = avaliacao["conceito"]
                    cor = avaliacao["cor"]
                    individual_evals[avaliador] = {"conceito": conceito, "cor": cor}

                result[dir_name][comp_name] = {
                    "scores": scores,
                    "individual_evaluations": individual_evals,
                }

        return result

    def generate_overview_sheet(self) -> pd.DataFrame:
        """Generate overview data with concepts and scores by year for all people"""
        overview_data = []

        for person in self.all_people:
            # Get concepts for all years
            concepts = self.get_concept_by_year(person)

            # Calculate average scores for each year
            person_row = {"Pessoa": person}

            for year in self.all_years:
                # Add concept for this year
                person_row[f"Conceito {year}"] = concepts.get(year, "n/a")

                # Calculate average score for this year
                if year in self.evaluations_by_person[person]:
                    total_score = 0
                    total_group_score = 0
                    count = 0

                    behavior_scores = self.get_behavior_scores(person, year)
                    for dir_name, behaviors in behavior_scores.items():
                        for comp_name, details in behaviors.items():
                            for avaliador, scores in details["scores"].items():
                                if avaliador == "%todos":  # Use the overall evaluation
                                    total_score += scores["score_colaborador"]
                                    total_group_score += scores["score_grupo"]
                                    count += 1

                    if count > 0:
                        person_row[f"Score {year}"] = total_score / count
                        person_row[f"Score Grupo {year}"] = total_group_score / count
                        person_row[f"Diferença {year}"] = (total_score / count) - (
                            total_group_score / count
                        )
                    else:
                        person_row[f"Score {year}"] = np.nan
                        person_row[f"Score Grupo {year}"] = np.nan
                        person_row[f"Diferença {year}"] = np.nan
                else:
                    person_row[f"Score {year}"] = np.nan
                    person_row[f"Score Grupo {year}"] = np.nan
                    person_row[f"Diferença {year}"] = np.nan

            overview_data.append(person_row)

        overview_df = pd.DataFrame(overview_data)

        # Calculate ranking for each year based on scores
        for year in self.all_years:
            score_col = f"Score {year}"
            if score_col in overview_df.columns:
                # Create ranking (1 is best, higher is worse)
                overview_df[f"Ranking {year}"] = overview_df[score_col].rank(
                    ascending=False, method="min", na_option="bottom"
                )

        return overview_df

    def generate_detailed_year_sheet(
        self, year: str
    ) -> Tuple[pd.DataFrame, Dict[str, List[str]]]:
        """Generate detailed comparison for a specific year"""
        if year not in self.year_criteria:
            return pd.DataFrame(), {}

        # First collect all behaviors for this year
        behavior_list = []
        for dir_name, behaviors in self.year_criteria[year].items():
            for behavior in behaviors:
                behavior_list.append((dir_name, behavior))

        # Sort by direcionador and then behavior
        behavior_list.sort()

        # Create a data frame with one row per person, columns for each behavior
        detailed_data = []
        # Keep track of comments for each person+behavior
        all_comments = defaultdict(list)

        for person in self.all_people:
            if year not in self.evaluations_by_person.get(person, {}):
                continue

            person_row = {"Pessoa": person}

            # Get overall concept for this year
            if (
                "data" in self.evaluations_by_person[person][year]
                and "conceito_ciclo_filho_descricao"
                in self.evaluations_by_person[person][year]["data"]
            ):
                person_row["Conceito"] = self.evaluations_by_person[person][year][
                    "data"
                ]["conceito_ciclo_filho_descricao"]
            else:
                person_row["Conceito"] = "n/a"

            # Get scores for each behavior
            behavior_scores = self.get_behavior_scores(person, year)

            for dir_name, behavior in behavior_list:
                col_name = f"{dir_name} - {behavior}"

                if (
                    dir_name in behavior_scores
                    and behavior in behavior_scores[dir_name]
                ):
                    details = behavior_scores[dir_name][behavior]

                    # Get overall score
                    if "%todos" in details["scores"]:
                        person_row[col_name] = details["scores"]["%todos"][
                            "score_colaborador"
                        ]
                        person_row[f"{col_name} (Grupo)"] = details["scores"]["%todos"][
                            "score_grupo"
                        ]
                        person_row[f"{col_name} (Diff)"] = details["scores"]["%todos"][
                            "difference"
                        ]
                    else:
                        person_row[col_name] = np.nan
                        person_row[f"{col_name} (Grupo)"] = np.nan
                        person_row[f"{col_name} (Diff)"] = np.nan

                    # Add self-evaluation and manager evaluation
                    for evaluator in ["autoavaliação", "gestor"]:
                        if evaluator in details["individual_evaluations"]:
                            eval_concept = details["individual_evaluations"][evaluator][
                                "conceito"
                            ]
                            person_row[f"{col_name} ({evaluator})"] = eval_concept

                            # Add to comments
                            comment_key = f"{person}|{col_name}"
                            all_comments[comment_key].append(
                                f"{evaluator}: {eval_concept}"
                            )
                        else:
                            person_row[f"{col_name} ({evaluator})"] = ""
                else:
                    person_row[col_name] = np.nan
                    person_row[f"{col_name} (Grupo)"] = np.nan
                    person_row[f"{col_name} (Diff)"] = np.nan
                    person_row[f"{col_name} (autoavaliação)"] = ""
                    person_row[f"{col_name} (gestor)"] = ""

            detailed_data.append(person_row)

        detailed_df = pd.DataFrame(detailed_data)

        # Calculate rankings for each behavior
        for dir_name, behavior in behavior_list:
            col_name = f"{dir_name} - {behavior}"
            if col_name in detailed_df.columns:
                # Create ranking (1 is best, higher is worse)
                detailed_df[f"{col_name} (Ranking)"] = detailed_df[col_name].rank(
                    ascending=False, method="min", na_option="bottom"
                )

        return detailed_df, all_comments

    def generate_comments_sheet(self) -> pd.DataFrame:
        """Generate a sheet with all comments and observations for each person and behavior"""
        comments_data = []

        # Process each person and year
        for person in self.all_people:
            for year in self.all_years:
                if year not in self.evaluations_by_person.get(person, {}):
                    continue

                # Get behavior scores and comments
                behavior_scores = self.get_behavior_scores(person, year)

                # Process each behavior
                for dir_name, behaviors in behavior_scores.items():
                    for behavior_name, details in behaviors.items():
                        # Get individual evaluations (comments)
                        for evaluator, eval_data in details[
                            "individual_evaluations"
                        ].items():
                            comments_data.append({
                                "Pessoa": person,
                                "Ano": year,
                                "Direcionador": dir_name,
                                "Comportamento": behavior_name,
                                "Avaliador": evaluator,
                                "Conceito": eval_data["conceito"],
                                "Score Colaborador": details["scores"]
                                .get("%todos", {})
                                .get("score_colaborador", np.nan),
                                "Score Grupo": details["scores"]
                                .get("%todos", {})
                                .get("score_grupo", np.nan),
                                "Diferença": details["scores"]
                                .get("%todos", {})
                                .get("difference", np.nan),
                            })

        # Create DataFrame
        comments_df = pd.DataFrame(comments_data)

        # Sort by person, year, direcionador, and comportamento
        if not comments_df.empty:
            comments_df = comments_df.sort_values([
                "Pessoa",
                "Ano",
                "Direcionador",
                "Comportamento",
            ])

        return comments_df

    def apply_excel_formatting(self, writer: pd.ExcelWriter):
        """Apply formatting to Excel workbook"""
        workbook = writer.book

        # Format for the overview sheet
        overview_sheet = writer.sheets["Visão Geral"]

        # Get dimensions directly from the sheet
        overview_df_cols = len(overview_sheet[1])
        overview_df_rows = overview_sheet.max_row

        # Apply header formatting
        header_fill = PatternFill(
            start_color="4472C4", end_color="4472C4", fill_type="solid"
        )
        header_font = Font(color="FFFFFF", bold=True)
        header_alignment = Alignment(
            horizontal="center", vertical="center", wrap_text=True
        )

        for col_idx in range(1, overview_df_cols + 1):
            cell = overview_sheet.cell(row=1, column=col_idx)
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = header_alignment

        # Apply conditional formatting for concept cells
        for row_idx in range(2, overview_df_rows + 1):
            for year in self.all_years:
                # Find the concept column for each year
                col_name = f"Conceito {year}"
                # Find the column index by checking header row
                col_idx = None
                for c_idx in range(1, overview_df_cols + 1):
                    if overview_sheet.cell(row=1, column=c_idx).value == col_name:
                        col_idx = c_idx
                        break

                if col_idx:
                    cell = overview_sheet.cell(row=row_idx, column=col_idx)
                    concept = cell.value

                    if concept in self.concept_colors:
                        cell.fill = PatternFill(
                            start_color=self.concept_colors[concept].replace("#", ""),
                            end_color=self.concept_colors[concept].replace("#", ""),
                            fill_type="solid",
                        )

        # Add conditional formatting for score differences
        for year in self.all_years:
            diff_col_idx = None
            for c_idx in range(1, overview_df_cols + 1):
                if (
                    overview_sheet.cell(row=1, column=c_idx).value
                    == f"Diferença {year}"
                ):
                    diff_col_idx = c_idx
                    break

            if diff_col_idx:
                # Add conditional formatting for score differences
                diff_range = f"{get_column_letter(diff_col_idx)}2:{get_column_letter(diff_col_idx)}{overview_df_rows}"

                # Add green-yellow-red color scale
                color_scale_rule = ColorScaleRule(
                    start_type="min",
                    start_color="FF0000",  # Red for negative
                    mid_type="num",
                    mid_value=0,
                    mid_color="FFFF00",  # Yellow for zero
                    end_type="max",
                    end_color="00FF00",  # Green for positive
                )
                overview_sheet.conditional_formatting.add(diff_range, color_scale_rule)

        # Add conditional formatting for rankings
        for year in self.all_years:
            rank_col_idx = None
            for c_idx in range(1, overview_df_cols + 1):
                if overview_sheet.cell(row=1, column=c_idx).value == f"Ranking {year}":
                    rank_col_idx = c_idx
                    break

            if rank_col_idx:
                # Add icons for top 3 ranking
                rank_range = f"{get_column_letter(rank_col_idx)}2:{get_column_letter(rank_col_idx)}{overview_df_rows}"

                # Add rule for top 3
                top3_rule = CellIsRule(
                    operator="lessThanOrEqual",
                    formula=["3"],
                    stopIfTrue=True,
                    fill=PatternFill(
                        start_color="92D050", end_color="92D050", fill_type="solid"
                    ),
                )
                overview_sheet.conditional_formatting.add(rank_range, top3_rule)

        # Apply similar formatting to year detail sheets
        for year in self.all_years:
            if f"Detalhes {year}" in writer.sheets:
                year_sheet = writer.sheets[f"Detalhes {year}"]

                # Get dimensions directly from the sheet
                year_df_cols = len(year_sheet[1])
                year_df_rows = year_sheet.max_row

                # Apply header formatting
                for col_idx in range(1, year_df_cols + 1):
                    cell = year_sheet.cell(row=1, column=col_idx)
                    cell.fill = header_fill
                    cell.font = header_font
                    cell.alignment = header_alignment

                # Apply conditional formatting for concept cells
                for row_idx in range(2, year_df_rows + 1):
                    # Find the column index for "Conceito"
                    concept_col_idx = None
                    for c_idx in range(1, year_df_cols + 1):
                        if year_sheet.cell(row=1, column=c_idx).value == "Conceito":
                            concept_col_idx = c_idx
                            break

                    if concept_col_idx:
                        cell = year_sheet.cell(row=row_idx, column=concept_col_idx)
                        concept = cell.value

                        if concept in self.concept_colors:
                            cell.fill = PatternFill(
                                start_color=self.concept_colors[concept].replace(
                                    "#", ""
                                ),
                                end_color=self.concept_colors[concept].replace("#", ""),
                                fill_type="solid",
                            )

                # Add conditional formatting for behavior scores
                for col_idx in range(1, year_df_cols + 1):
                    col_header = year_sheet.cell(row=1, column=col_idx).value
                    if (
                        col_header
                        and isinstance(col_header, str)
                        and " - " in col_header
                        and not any(
                            suffix in col_header
                            for suffix in [
                                "(Grupo)",
                                "(Diff)",
                                "(Ranking)",
                                "(autoavaliação)",
                                "(gestor)",
                            ]
                        )
                    ):
                        # This is a behavior score column
                        behavior_range = f"{get_column_letter(col_idx)}2:{get_column_letter(col_idx)}{year_df_rows}"

                        # Add color scale for scores (1-5 scale)
                        color_scale_rule = ColorScaleRule(
                            start_type="num",
                            start_value=1,
                            start_color="FF6961",  # Red for low scores
                            mid_type="num",
                            mid_value=3,
                            mid_color="FFFF99",  # Yellow for middle scores
                            end_type="num",
                            end_value=5,
                            end_color="77DD77",  # Green for high scores
                        )
                        year_sheet.conditional_formatting.add(
                            behavior_range, color_scale_rule
                        )

                        # Find corresponding difference column
                        diff_col_idx = None
                        diff_col_header = f"{col_header} (Diff)"
                        for c_idx in range(1, year_df_cols + 1):
                            if (
                                year_sheet.cell(row=1, column=c_idx).value
                                == diff_col_header
                            ):
                                diff_col_idx = c_idx
                                break

                        if diff_col_idx:
                            # Add conditional formatting for differences
                            diff_range = f"{get_column_letter(diff_col_idx)}2:{get_column_letter(diff_col_idx)}{year_df_rows}"

                            # Add green-yellow-red color scale for differences
                            diff_scale_rule = ColorScaleRule(
                                start_type="min",
                                start_color="FF0000",  # Red for negative
                                mid_type="num",
                                mid_value=0,
                                mid_color="FFFF00",  # Yellow for zero
                                end_type="max",
                                end_color="00FF00",  # Green for positive
                            )
                            year_sheet.conditional_formatting.add(
                                diff_range, diff_scale_rule
                            )

                # Set column widths
                for col_idx in range(1, year_df_cols + 1):
                    year_sheet.column_dimensions[get_column_letter(col_idx)].width = 20

                # Add borders to all cells
                thin_border = Border(
                    left=Side(style="thin"),
                    right=Side(style="thin"),
                    top=Side(style="thin"),
                    bottom=Side(style="thin"),
                )

                for row_idx in range(1, year_df_rows + 1):
                    for col_idx in range(1, year_df_cols + 1):
                        cell = year_sheet.cell(row=row_idx, column=col_idx)
                        cell.border = thin_border

        # Add formatting to Trends sheet
        if "Tendências" in writer.sheets:
            trends_sheet = writer.sheets["Tendências"]

            # Get dimensions
            trends_df_cols = len(trends_sheet[1])
            trends_df_rows = trends_sheet.max_row

            # Apply header formatting
            for col_idx in range(1, trends_df_cols + 1):
                cell = trends_sheet.cell(row=1, column=col_idx)
                cell.fill = header_fill
                cell.font = header_font
                cell.alignment = header_alignment

            # Add formatting for variation columns
            for col_idx in range(1, trends_df_cols + 1):
                col_header = trends_sheet.cell(row=1, column=col_idx).value
                if (
                    col_header
                    and isinstance(col_header, str)
                    and col_header.startswith("Variação")
                ):
                    # Apply color scale to variation
                    var_range = f"{get_column_letter(col_idx)}2:{get_column_letter(col_idx)}{trends_df_rows}"

                    # Add green-yellow-red color scale
                    color_scale_rule = ColorScaleRule(
                        start_type="min",
                        start_color="FF0000",  # Red for negative
                        mid_type="num",
                        mid_value=0,
                        mid_color="FFFF00",  # Yellow for zero
                        end_type="max",
                        end_color="00FF00",  # Green for positive
                    )
                    trends_sheet.conditional_formatting.add(var_range, color_scale_rule)

            # Add concept formatting
            for row_idx in range(2, trends_df_rows + 1):
                # Find concept columns
                for concept_header in ["Conceito Primeiro Ano", "Conceito Último Ano"]:
                    concept_col_idx = None
                    for c_idx in range(1, trends_df_cols + 1):
                        if (
                            trends_sheet.cell(row=1, column=c_idx).value
                            == concept_header
                        ):
                            concept_col_idx = c_idx
                            break

                    if concept_col_idx:
                        cell = trends_sheet.cell(row=row_idx, column=concept_col_idx)
                        concept = cell.value

                        if concept in self.concept_colors:
                            cell.fill = PatternFill(
                                start_color=self.concept_colors[concept].replace(
                                    "#", ""
                                ),
                                end_color=self.concept_colors[concept].replace("#", ""),
                                fill_type="solid",
                            )

            # Set column widths
            for col_idx in range(1, trends_df_cols + 1):
                trends_sheet.column_dimensions[get_column_letter(col_idx)].width = 20

        # Add formatting to Comments sheet
        if "Comentários" in writer.sheets:
            comments_sheet = writer.sheets["Comentários"]

            # Get dimensions
            comments_df_cols = len(comments_sheet[1])
            comments_df_rows = comments_sheet.max_row

            # Apply header formatting
            for col_idx in range(1, comments_df_cols + 1):
                cell = comments_sheet.cell(row=1, column=col_idx)
                cell.fill = header_fill
                cell.font = header_font
                cell.alignment = header_alignment

            # Set column widths
            for col_idx in range(1, comments_df_cols + 1):
                comments_sheet.column_dimensions[get_column_letter(col_idx)].width = 20

            # Add alternating row colors for readability
            for row_idx in range(2, comments_df_rows + 1):
                if row_idx % 2 == 0:  # Even rows
                    for col_idx in range(1, comments_df_cols + 1):
                        cell = comments_sheet.cell(row=row_idx, column=col_idx)
                        cell.fill = PatternFill(
                            start_color="F2F2F2", end_color="F2F2F2", fill_type="solid"
                        )

    def generate_trends_sheet(self) -> pd.DataFrame:
        """Generate a sheet with performance trends over time for each person"""
        trends_data = []

        # For each person, calculate their performance trend over years
        for person in self.all_people:
            person_years = [
                year
                for year in self.all_years
                if year in self.evaluations_by_person.get(person, {})
            ]

            if not person_years:
                continue

            # Get concepts for available years
            concepts = self.get_concept_by_year(person)

            # Calculate performance metrics for each year
            year_metrics = {}
            for year in person_years:
                # Calculate average scores
                total_score = 0
                total_group_score = 0
                count = 0

                # Track scores by direcionador
                direcionador_scores = defaultdict(lambda: {"score": 0, "count": 0})

                behavior_scores = self.get_behavior_scores(person, year)
                for dir_name, behaviors in behavior_scores.items():
                    for comp_name, details in behaviors.items():
                        for avaliador, scores in details["scores"].items():
                            if avaliador == "%todos":  # Use the overall evaluation
                                total_score += scores["score_colaborador"]
                                total_group_score += scores["score_grupo"]
                                count += 1

                                # Track by direcionador
                                direcionador_scores[dir_name]["score"] += scores[
                                    "score_colaborador"
                                ]
                                direcionador_scores[dir_name]["count"] += 1

                # Overall average
                year_metrics[year] = {
                    "avg_score": total_score / count if count > 0 else np.nan,
                    "avg_group": total_group_score / count if count > 0 else np.nan,
                    "direcionador_scores": {
                        dir_name: data["score"] / data["count"]
                        if data["count"] > 0
                        else np.nan
                        for dir_name, data in direcionador_scores.items()
                    },
                }

            # Get all unique direcionadores across all years for this person
            all_direcionadores = set()
            for year in person_years:
                all_direcionadores.update(
                    year_metrics[year]["direcionador_scores"].keys()
                )

            # Calculate trend for overall score
            if len(person_years) >= 2:
                first_year = min(person_years)
                last_year = max(person_years)

                overall_trend = (
                    year_metrics[last_year]["avg_score"]
                    - year_metrics[first_year]["avg_score"]
                )

                # Calculate trends for each direcionador
                direcionador_trends = {}
                for dir_name in all_direcionadores:
                    # Check if direcionador exists in both first and last year
                    if (
                        dir_name in year_metrics[first_year]["direcionador_scores"]
                        and dir_name in year_metrics[last_year]["direcionador_scores"]
                    ):
                        first_score = year_metrics[first_year]["direcionador_scores"][
                            dir_name
                        ]
                        last_score = year_metrics[last_year]["direcionador_scores"][
                            dir_name
                        ]

                        if not np.isnan(first_score) and not np.isnan(last_score):
                            direcionador_trends[dir_name] = last_score - first_score
            else:
                overall_trend = np.nan
                direcionador_trends = {}

            # Create row for person
            person_row = {
                "Pessoa": person,
                "Primeiro Ano": min(person_years) if person_years else None,
                "Último Ano": max(person_years) if person_years else None,
                "Conceito Primeiro Ano": concepts.get(min(person_years), "n/a")
                if person_years
                else "n/a",
                "Conceito Último Ano": concepts.get(max(person_years), "n/a")
                if person_years
                else "n/a",
                "Variação no Score": overall_trend,
            }

            # Add scores for each year
            for year in self.all_years:
                if year in year_metrics:
                    person_row[f"Score {year}"] = year_metrics[year]["avg_score"]
                    person_row[f"Score Grupo {year}"] = year_metrics[year]["avg_group"]
                    person_row[f"Diferença {year}"] = (
                        year_metrics[year]["avg_score"]
                        - year_metrics[year]["avg_group"]
                    )
                else:
                    person_row[f"Score {year}"] = np.nan
                    person_row[f"Score Grupo {year}"] = np.nan
                    person_row[f"Diferença {year}"] = np.nan

            # Add trends for each direcionador
            for dir_name in sorted(all_direcionadores):
                trend_value = direcionador_trends.get(dir_name, np.nan)
                person_row[f"Variação {dir_name}"] = trend_value

            trends_data.append(person_row)

        # Create DataFrame
        trends_df = pd.DataFrame(trends_data)

        # Sort by variation in score (descending)
        if not trends_df.empty and "Variação no Score" in trends_df.columns:
            trends_df = trends_df.sort_values("Variação no Score", ascending=False)

        return trends_df

    def add_embedded_charts(self, writer: pd.ExcelWriter):
        """Add embedded charts to the workbook"""
        workbook = writer.book

        # Add radar charts to the overview sheet
        self._add_radar_charts(workbook)

        # Add trend charts
        self._add_trend_charts(workbook)

        # Add distribution charts
        self._add_distribution_charts(workbook)

    def _add_radar_charts(self, workbook: Workbook):
        """Add radar charts for top performers"""
        if "Visão Geral" not in workbook.sheetnames:
            return

        sheet = workbook["Visão Geral"]

        # Get data
        last_year = max(self.all_years) if self.all_years else None
        if not last_year:
            return

        # Find columns containing last year data
        score_col = None
        for col in range(1, sheet.max_column + 1):
            if sheet.cell(row=1, column=col).value == f"Score {last_year}":
                score_col = col
                break

        if not score_col:
            return

        # Find top 3 performers in the last year
        top_performers = []
        scores = {}

        for row in range(2, sheet.max_row + 1):
            person = sheet.cell(row=row, column=1).value
            score = sheet.cell(row=row, column=score_col).value

            if score is not None and not isinstance(score, str):
                scores[person] = score

        # Sort by score and get top 3
        top_3 = sorted(scores.items(), key=lambda x: x[1], reverse=True)[:3]
        top_performers = [p[0] for p in top_3]

        if not top_performers:
            return

        # Create radar chart for top performers
        if f"Detalhes {last_year}" not in workbook.sheetnames:
            return

        detail_sheet = workbook[f"Detalhes {last_year}"]

        # Find all behavior columns (without the suffixes)
        behavior_cols = {}
        for col in range(1, detail_sheet.max_column + 1):
            header = detail_sheet.cell(row=1, column=col).value
            if (
                header
                and isinstance(header, str)
                and " - " in header
                and not any(
                    suffix in header
                    for suffix in [
                        "(Grupo)",
                        "(Diff)",
                        "(Ranking)",
                        "(autoavaliação)",
                        "(gestor)",
                    ]
                )
            ):
                behavior_cols[header] = col

        if not behavior_cols:
            return

        # Create a radar chart
        chart = RadarChart()
        chart.type = "filled"
        chart.style = 26
        chart.title = f"Top Performers {last_year} - Comportamentos"

        # Add data for each top performer
        for person in top_performers:
            # Find the row for this person
            person_row = None
            for row in range(2, detail_sheet.max_row + 1):
                if detail_sheet.cell(row=row, column=1).value == person:
                    person_row = row
                    break

            if not person_row:
                continue

            # Create data for this person
            data = []
            for behavior, col in behavior_cols.items():
                value = detail_sheet.cell(row=person_row, column=col).value
                if value is None or isinstance(value, str):
                    value = 0
                data.append(value)

            # If we have data, add it to the chart
            if data:
                # Add behavior names as categories
                cats = Reference(detail_sheet, min_col=1, min_row=2, max_row=2)

                # Add data series for this person
                values = Reference(
                    detail_sheet,
                    min_col=2,
                    min_row=person_row,
                    max_col=len(data) + 1,
                    max_row=person_row,
                )
                series = Series(values, title=person)
                chart.series.append(series)

        # Size and position the chart
        chart.width = 15
        chart.height = 10

        # Add the chart to the overview sheet
        detail_sheet.add_chart(chart, "A" + str(detail_sheet.max_row + 5))

    def _add_trend_charts(self, workbook: Workbook):
        """Add trend charts to the workbook"""
        if "Tendências" not in workbook.sheetnames:
            return

        sheet = workbook["Tendências"]

        # Create a line chart for score evolution
        chart = LineChart()
        chart.title = "Evolução de Performance"
        chart.style = 12
        chart.x_axis.title = "Ano"
        chart.y_axis.title = "Score Médio"
        chart.width = 20
        chart.height = 10

        # Add X-axis categories (years)
        years = sorted(self.all_years)
        cats = Reference(sheet, min_col=1, min_row=1, max_row=1, max_col=len(years))
        chart.set_categories(cats)

        # Find columns with yearly scores
        score_cols = {}
        for year in years:
            for col in range(1, sheet.max_column + 1):
                if sheet.cell(row=1, column=col).value == f"Score {year}":
                    score_cols[year] = col
                    break

        # Add a data series for each person
        for row in range(
            2, min(sheet.max_row + 1, 7)
        ):  # Limit to first 5 people for readability
            person = sheet.cell(row=row, column=1).value

            # Create data series
            values = []
            for year in years:
                if year in score_cols:
                    col = score_cols[year]
                    values.append(sheet.cell(row=row, column=col).value)
                else:
                    values.append(None)

            # Add to chart if we have data
            if any(v is not None for v in values):
                data = Reference(
                    sheet, min_col=1, min_row=row, max_row=row, max_col=len(values)
                )
                series = Series(data, title=person)
                chart.series.append(series)

        # Add chart to the trends sheet
        sheet.add_chart(chart, "A" + str(sheet.max_row + 5))

    def _add_distribution_charts(self, workbook: Workbook):
        """Add distribution charts for concepts"""
        if "Visão Geral" not in workbook.sheetnames:
            return

        sheet = workbook["Visão Geral"]

        # Create a pie chart for concept distribution in the latest year
        last_year = max(self.all_years) if self.all_years else None
        if not last_year:
            return

        # Find the column with concepts for the last year
        concept_col = None
        for col in range(1, sheet.max_column + 1):
            if sheet.cell(row=1, column=col).value == f"Conceito {last_year}":
                concept_col = col
                break

        if not concept_col:
            return

        # Count concepts
        concepts = defaultdict(int)
        for row in range(2, sheet.max_row + 1):
            concept = sheet.cell(row=row, column=concept_col).value
            if concept:
                concepts[concept] += 1

        if not concepts:
            return

        # Create a temporary sheet for the chart data
        if "ChartData" not in workbook.sheetnames:
            workbook.create_sheet("ChartData")
            workbook["ChartData"].sheet_state = "hidden"  # Hide this sheet

        chart_sheet = workbook["ChartData"]

        # Add concept distribution data
        chart_sheet.cell(row=1, column=1).value = "Conceito"
        chart_sheet.cell(row=1, column=2).value = "Quantidade"

        for i, (concept, count) in enumerate(concepts.items(), 2):
            chart_sheet.cell(row=i, column=1).value = concept
            chart_sheet.cell(row=i, column=2).value = count

        # Create pie chart
        pie = PieChart()
        pie.title = f"Distribuição de Conceitos - {last_year}"

        # Add data
        labels = Reference(chart_sheet, min_col=1, min_row=2, max_row=1 + len(concepts))
        data = Reference(chart_sheet, min_col=2, min_row=1, max_row=1 + len(concepts))
        pie.add_data(data, titles_from_data=True)
        pie.set_categories(labels)

        # Add to the overview sheet
        sheet.add_chart(pie, "J2")

    def generate_gap_analysis_sheet(self) -> pd.DataFrame:
        """Generate a sheet with gap analysis and recommendations for each person"""
        gap_data = []

        # For each person, identify top 3 gaps and strengths
        for person in self.all_people:
            last_year = max(self.all_years) if self.all_years else None
            if not last_year or last_year not in self.evaluations_by_person.get(
                person, {}
            ):
                continue

            # Get behavior scores for the last year
            behavior_scores = self.get_behavior_scores(person, last_year)

            # Collect all behaviors and their scores
            all_behaviors = []
            for dir_name, behaviors in behavior_scores.items():
                for comp_name, details in behaviors.items():
                    for avaliador, scores in details["scores"].items():
                        if avaliador == "%todos":  # Use the overall evaluation
                            score = scores["score_colaborador"]
                            group_score = scores["score_grupo"]
                            difference = scores["difference"]

                            all_behaviors.append({
                                "direcionador": dir_name,
                                "comportamento": comp_name,
                                "score": score,
                                "group_score": group_score,
                                "difference": difference,
                            })

            # Sort by difference (ascending for gaps, descending for strengths)
            all_behaviors.sort(key=lambda x: x["difference"])

            # Top 3 gaps (worst differences)
            top_gaps = all_behaviors[:3] if len(all_behaviors) >= 3 else all_behaviors

            # Top 3 strengths (best differences)
            top_strengths = (
                all_behaviors[-3:] if len(all_behaviors) >= 3 else all_behaviors
            )
            top_strengths.reverse()  # Reverse to get descending order

            # Calculate average score
            avg_score = (
                sum(b["score"] for b in all_behaviors) / len(all_behaviors)
                if all_behaviors
                else 0
            )
            avg_group = (
                sum(b["group_score"] for b in all_behaviors) / len(all_behaviors)
                if all_behaviors
                else 0
            )

            # Generate recommendations based on gaps
            recommendations = self._generate_recommendations(top_gaps)

            # Create row for this person
            gap_row = {
                "Pessoa": person,
                "Conceito": self.get_concept_by_year(person).get(last_year, "n/a"),
                "Score Médio": avg_score,
                "Score Grupo": avg_group,
                "Diferença": avg_score - avg_group,
                "Gap Principal": f"{top_gaps[0]['direcionador']} - {top_gaps[0]['comportamento']}"
                if top_gaps
                else "",
                "Score Gap Principal": top_gaps[0]["score"] if top_gaps else np.nan,
                "Diferença Gap Principal": top_gaps[0]["difference"]
                if top_gaps
                else np.nan,
                "Gap Secundário": f"{top_gaps[1]['direcionador']} - {top_gaps[1]['comportamento']}"
                if len(top_gaps) > 1
                else "",
                "Score Gap Secundário": top_gaps[1]["score"]
                if len(top_gaps) > 1
                else np.nan,
                "Diferença Gap Secundário": top_gaps[1]["difference"]
                if len(top_gaps) > 1
                else np.nan,
                "Gap Terciário": f"{top_gaps[2]['direcionador']} - {top_gaps[2]['comportamento']}"
                if len(top_gaps) > 2
                else "",
                "Score Gap Terciário": top_gaps[2]["score"]
                if len(top_gaps) > 2
                else np.nan,
                "Diferença Gap Terciário": top_gaps[2]["difference"]
                if len(top_gaps) > 2
                else np.nan,
                "Força Principal": f"{top_strengths[0]['direcionador']} - {top_strengths[0]['comportamento']}"
                if top_strengths
                else "",
                "Score Força Principal": top_strengths[0]["score"]
                if top_strengths
                else np.nan,
                "Diferença Força Principal": top_strengths[0]["difference"]
                if top_strengths
                else np.nan,
                "Força Secundária": f"{top_strengths[1]['direcionador']} - {top_strengths[1]['comportamento']}"
                if len(top_strengths) > 1
                else "",
                "Score Força Secundária": top_strengths[1]["score"]
                if len(top_strengths) > 1
                else np.nan,
                "Diferença Força Secundária": top_strengths[1]["difference"]
                if len(top_strengths) > 1
                else np.nan,
                "Força Terciária": f"{top_strengths[2]['direcionador']} - {top_strengths[2]['comportamento']}"
                if len(top_strengths) > 2
                else "",
                "Score Força Terciária": top_strengths[2]["score"]
                if len(top_strengths) > 2
                else np.nan,
                "Diferença Força Terciária": top_strengths[2]["difference"]
                if len(top_strengths) > 2
                else np.nan,
                "Recomendações": recommendations,
            }

            gap_data.append(gap_row)

        # Create DataFrame
        gap_df = pd.DataFrame(gap_data)

        # Sort by average score difference
        if not gap_df.empty and "Diferença" in gap_df.columns:
            gap_df = gap_df.sort_values("Diferença")

        return gap_df

    def _generate_recommendations(self, gaps: List[Dict]) -> str:
        """Generate recommendations based on identified gaps"""
        if not gaps:
            return "Não foram identificados gaps significativos."

        recommendations = []

        # Dictionary of recommendations by direcionador
        rec_by_dir = {
            "A gente trabalha para o cliente": [
                "Participar de workshops de design thinking centrado no cliente",
                "Implementar pesquisas de satisfação para suas entregas",
                "Acompanhar de perto a jornada do cliente para entender pontos de atrito",
                "Conversar diretamente com clientes para entender suas necessidades reais",
            ],
            "Performance que transforma": [
                "Definir métricas claras para resultados sustentáveis",
                "Participar de projetos inovadores fora da zona de conforto",
                "Criar plano de desenvolvimento específico para mentalidade de dono",
                "Estruturar análises de impacto para decisões importantes",
            ],
            "Liderança mobilizadora": [
                "Participar de um programa de mentoria para desenvolver habilidades de liderança",
                "Buscar feedback contínuo da equipe sobre seu estilo de comunicação",
                "Implementar práticas de reconhecimento e celebração de conquistas",
                "Desenvolver um plano de desenvolvimento para cada membro da equipe",
            ],
            "Pensamento digital": [
                "Participar de cursos sobre transformação digital e novas tecnologias",
                "Propor soluções que integrem tecnologia para resolver problemas",
                "Implementar processos ágeis em seu trabalho diário",
                "Explorar ferramentas digitais que possam aumentar a eficiência",
            ],
            "Dono do negócio": [
                "Desenvolver visão estratégica através de benchmarks com concorrentes",
                "Implementar acompanhamento sistemático de resultados",
                "Criar planos de ação para mitigar riscos identificados",
                "Participar de decisões estratégicas além do escopo atual",
            ],
        }

        # Generic recommendations if direcionador not found
        generic_recs = [
            "Buscar mentoria específica para desenvolver esta competência",
            "Criar plano de desenvolvimento pessoal com foco nesta área",
            "Participar de treinamentos específicos neste tema",
            "Buscar feedback estruturado sobre este comportamento",
        ]

        # Add specific recommendations for each gap
        for i, gap in enumerate(gaps):
            dir_name = gap["direcionador"]
            comp_name = gap["comportamento"]

            if i == 0:
                recommendations.append(
                    f'**Prioridade 1:** Para melhorar em "{comp_name}":'
                )
            else:
                recommendations.append(
                    f'**Prioridade {i + 1}:** Para melhorar em "{comp_name}":'
                )

            # Get recommendations for this direcionador
            dir_key = next((k for k in rec_by_dir.keys() if k in dir_name), None)
            dir_recs = rec_by_dir.get(dir_key, generic_recs)

            # Add 1-2 recommendations
            for j in range(min(2, len(dir_recs))):
                idx = (i + j) % len(dir_recs)  # Rotate through recommendations
                recommendations.append(f"  • {dir_recs[idx]}")

            recommendations.append("")  # Add empty line

        # Add timeline suggestion
        recommendations.append("**Plano de Ação Sugerido:**")
        recommendations.append(
            "• Curto prazo (3 meses): Foco no gap principal com feedback frequente"
        )
        recommendations.append(
            "• Médio prazo (6 meses): Desenvolvimento dos gaps secundários"
        )
        recommendations.append(
            "• Longo prazo (12 meses): Consolidação e mentoria para outros"
        )

        return "\n".join(recommendations)

    def generate_comparative_spreadsheet(
        self, output_file: str = "avaliacao_comparativa.xlsx", export_pptx: bool = False
    ):
        """Generate comprehensive comparative spreadsheet with all data"""
        # Create Excel writer
        with pd.ExcelWriter(output_file, engine="openpyxl") as writer:
            # Generate overview sheet with concepts and scores by year
            print("Gerando aba de Visão Geral...")
            overview_df = self.generate_overview_sheet()
            overview_df.to_excel(writer, sheet_name="Visão Geral", index=False)

            # Generate trends sheet
            print("Gerando aba de Tendências...")
            trends_df = self.generate_trends_sheet()
            if not trends_df.empty:
                trends_df.to_excel(writer, sheet_name="Tendências", index=False)

            # Generate gap analysis sheet
            print("Gerando aba de Análise de Gaps...")
            gap_df = self.generate_gap_analysis_sheet()
            if not gap_df.empty:
                gap_df.to_excel(writer, sheet_name="Análise de Gaps", index=False)

            # Generate detailed sheets for each year
            for year in self.all_years:
                print(f"Gerando aba de Detalhes para {year}...")
                detailed_df, comments = self.generate_detailed_year_sheet(year)
                if not detailed_df.empty:
                    detailed_df.to_excel(
                        writer, sheet_name=f"Detalhes {year}", index=False
                    )

            # Generate comments sheet
            print("Gerando aba de Comentários...")
            comments_df = self.generate_comments_sheet()
            if not comments_df.empty:
                comments_df.to_excel(writer, sheet_name="Comentários", index=False)

            # Apply Excel formatting
            print("Aplicando formatação...")
            self.apply_excel_formatting(writer)

            # Add embedded charts
            print("Adicionando gráficos...")
            self.add_embedded_charts(writer)

        print(f"Planilha gerada: {output_file}")

        # Export to PowerPoint if requested
        if export_pptx and PPTX_AVAILABLE:
            print("Exportando para PowerPoint...")
            pptx_file = self.export_to_powerpoint(output_file)
            if pptx_file:
                print(f"Apresentação PowerPoint gerada: {pptx_file}")

        return output_file

    def export_to_powerpoint(
        self, excel_file: str, output_file: str = None
    ) -> Optional[str]:
        """Export key insights to PowerPoint presentation"""
        if not PPTX_AVAILABLE:
            print(
                "python-pptx não está instalado. A exportação para PowerPoint não está disponível."
            )
            return None

        if output_file is None:
            output_file = excel_file.replace(".xlsx", ".pptx")

        # Create presentation
        prs = Presentation()

        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Análise Comparativa 360°"
        subtitle.text = f"Gerado em {pd.Timestamp.now().strftime('%d/%m/%Y')}"

        # Load Excel file to get data
        wb = load_workbook(excel_file)

        # Add overview slide
        if "Visão Geral" in wb.sheetnames:
            self._add_overview_slide(prs, wb["Visão Geral"])

        # Add top performers slide
        if "Visão Geral" in wb.sheetnames:
            self._add_top_performers_slide(prs, wb["Visão Geral"])

        # Add gap analysis slide
        if "Análise de Gaps" in wb.sheetnames:
            self._add_gap_analysis_slide(prs, wb["Análise de Gaps"])

        # Add trends slide
        if "Tendências" in wb.sheetnames:
            self._add_trends_slide(prs, wb["Tendências"])

        # Save the presentation
        prs.save(output_file)
        print(f"Apresentação PowerPoint gerada: {output_file}")

        return output_file

    def _add_overview_slide(self, prs: "Presentation", sheet: "Worksheet"):
        """Add overview slide with key metrics"""
        # Add section title slide
        section_slide_layout = prs.slide_layouts[2]  # Section header layout
        slide = prs.slides.add_slide(section_slide_layout)
        title = slide.shapes.title
        title.text = "Visão Geral da Avaliação"

        # Add content slide
        content_slide_layout = prs.slide_layouts[5]  # Title and content layout
        slide = prs.slides.add_slide(content_slide_layout)
        title = slide.shapes.title
        title.text = "Resumo dos Resultados"

        # Get last year from sheet
        last_year = None
        for col in range(1, sheet.max_column + 1):
            header = sheet.cell(row=1, column=col).value
            if header and isinstance(header, str) and header.startswith("Conceito "):
                year = header.replace("Conceito ", "")
                if last_year is None or year > last_year:
                    last_year = year

        if not last_year:
            return

        # Find concept distribution
        concepts = defaultdict(int)
        concept_col = None
        for col in range(1, sheet.max_column + 1):
            if sheet.cell(row=1, column=col).value == f"Conceito {last_year}":
                concept_col = col
                break

        if concept_col:
            for row in range(2, sheet.max_row + 1):
                concept = sheet.cell(row=row, column=concept_col).value
                if concept:
                    concepts[concept] += 1

        # Create a table with concept distribution
        if concepts:
            # Add some text
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.text = f"Distribuição de Conceitos em {last_year}:"

            # Create a table below
            rows = len(concepts) + 1  # Header + data
            cols = 2  # Concept, Count
            left = Inches(2)
            top = Inches(2)
            width = Inches(6)
            height = Inches(0.5 * rows)

            table = slide.shapes.add_table(rows, cols, left, top, width, height).table

            # Add headers
            table.cell(0, 0).text = "Conceito"
            table.cell(0, 1).text = "Quantidade"

            # Add data
            for i, (concept, count) in enumerate(concepts.items(), 1):
                table.cell(i, 0).text = concept
                table.cell(i, 1).text = str(count)

                # Color the concept cell based on concept
                if concept in self.concept_colors:
                    cell = table.cell(i, 0)
                    fill = cell.fill
                    fill.solid()
                    rgb = self.concept_colors[concept].replace("#", "")
                    fill.fore_color.rgb = RGBColor(
                        int(rgb[0:2], 16), int(rgb[2:4], 16), int(rgb[4:6], 16)
                    )

    def _add_top_performers_slide(self, prs: "Presentation", sheet: "Worksheet"):
        """Add slide with top performers"""
        # Find last year and score column
        last_year = None
        score_col = None

        for col in range(1, sheet.max_column + 1):
            header = sheet.cell(row=1, column=col).value
            if header and isinstance(header, str) and header.startswith("Score "):
                year = header.replace("Score ", "")
                if last_year is None or year > last_year:
                    last_year = year
                    score_col = col

        if not last_year or not score_col:
            return

        # Get top 5 performers
        scores = []
        for row in range(2, sheet.max_row + 1):
            person = sheet.cell(row=row, column=1).value
            score = sheet.cell(row=row, column=score_col).value

            if score is not None and not isinstance(score, str):
                scores.append((person, score))

        # Sort by score
        scores.sort(key=lambda x: x[1], reverse=True)
        top_5 = scores[:5]

        if not top_5:
            return

        # Add slide
        slide_layout = prs.slide_layouts[5]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f"Top Performers em {last_year}"

        # Create a table
        rows = len(top_5) + 1  # Header + data
        cols = 3  # Rank, Person, Score
        left = Inches(1.5)
        top = Inches(2)
        width = Inches(7)
        height = Inches(0.5 * rows)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Add headers
        table.cell(0, 0).text = "Posição"
        table.cell(0, 1).text = "Pessoa"
        table.cell(0, 2).text = "Score"

        # Add data
        for i, (person, score) in enumerate(top_5, 1):
            table.cell(i, 0).text = str(i)
            table.cell(i, 1).text = person
            table.cell(i, 2).text = f"{score:.2f}"

            # Add gold/silver/bronze coloring
            if i <= 3:
                cell = table.cell(i, 0)
                fill = cell.fill
                fill.solid()

                if i == 1:
                    fill.fore_color.rgb = RGBColor(255, 215, 0)  # Gold
                elif i == 2:
                    fill.fore_color.rgb = RGBColor(192, 192, 192)  # Silver
                elif i == 3:
                    fill.fore_color.rgb = RGBColor(205, 127, 50)  # Bronze

    def _add_gap_analysis_slide(self, prs: "Presentation", sheet: "Worksheet"):
        """Add slide with gap analysis"""
        # Add slide
        slide_layout = prs.slide_layouts[5]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Análise de Gaps e Recomendações"

        # Add text frame
        content = slide.placeholders[1]
        tf = content.text_frame

        # Find people with largest gaps
        gaps = []
        for row in range(2, min(sheet.max_row + 1, 7)):  # Limit to top 5
            person = sheet.cell(row=row, column=1).value

            # Find gap columns
            gap_col = None
            for col in range(1, sheet.max_column + 1):
                if sheet.cell(row=1, column=col).value == "Gap Principal":
                    gap_col = col
                    break

            if gap_col and person:
                gap = sheet.cell(row=row, column=gap_col).value
                gaps.append((person, gap))

        if not gaps:
            tf.text = "Não foram encontrados gaps significativos."
            return

        # Add common areas for development
        p = tf.add_paragraph()
        p.text = "Principais Áreas para Desenvolvimento:"
        p.level = 0

        for person, gap in gaps:
            p = tf.add_paragraph()
            p.text = f"{person}: {gap}"
            p.level = 1

        # Add some recommendations
        p = tf.add_paragraph()
        p.text = "\nRecomendações Gerais:"
        p.level = 0

        recommendations = [
            "Implementar programa de desenvolvimento focado nas competências de cliente",
            "Estabelecer mentoria para liderança mobilizadora",
            "Criar workshops para pensamento digital e inovação",
            "Desenvolver planos de ação individuais baseados nos gaps identificados",
        ]

        for rec in recommendations:
            p = tf.add_paragraph()
            p.text = rec
            p.level = 1

    def _add_trends_slide(self, prs: "Presentation", sheet: "Worksheet"):
        """Add slide with performance trends"""
        # Add slide
        slide_layout = prs.slide_layouts[5]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Tendências de Performance"

        # Add text frame
        content = slide.placeholders[1]
        tf = content.text_frame

        # Find people with largest improvements
        improvements = []
        variation_col = None

        for col in range(1, sheet.max_column + 1):
            if sheet.cell(row=1, column=col).value == "Variação no Score":
                variation_col = col
                break

        if variation_col:
            for row in range(2, min(sheet.max_row + 1, 7)):  # Limit to top 5
                person = sheet.cell(row=row, column=1).value
                variation = sheet.cell(row=row, column=variation_col).value

                if person and variation is not None and not isinstance(variation, str):
                    improvements.append((person, variation))

            # Sort by variation
            improvements.sort(key=lambda x: x[1], reverse=True)

        if not improvements:
            tf.text = "Não foram encontradas tendências significativas."
            return

        # Add top improvers
        p = tf.add_paragraph()
        p.text = "Maiores Evoluções:"
        p.level = 0

        for person, variation in improvements[:3]:  # Top 3
            p = tf.add_paragraph()
            p.text = f"{person}: +{variation:.2f} pontos"
            p.level = 1

        # Add those who need attention
        p = tf.add_paragraph()
        p.text = "\nNecessitam Atenção:"
        p.level = 0

        for person, variation in sorted(improvements, key=lambda x: x[1])[
            :3
        ]:  # Bottom 3
            p = tf.add_paragraph()
            if variation < 0:
                p.text = f"{person}: {variation:.2f} pontos"
            else:
                p.text = f"{person}: +{variation:.2f} pontos"
            p.level = 1


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(
        description="Generate comparative spreadsheet of evaluations"
    )
    parser.add_argument("data_dir", help="Directory containing evaluation data")
    parser.add_argument(
        "--output", "-o", default="avaliacao_comparativa.xlsx", help="Output Excel file"
    )
    parser.add_argument(
        "--export-pptx",
        "-p",
        action="store_true",
        help="Export to PowerPoint presentation",
    )

    args = parser.parse_args()

    generator = ComparativeSpreadsheetGenerator(args.data_dir)
    output_file = generator.generate_comparative_spreadsheet(
        args.output, args.export_pptx
    )

    print(f"Planilha comparativa gerada: {output_file}")
